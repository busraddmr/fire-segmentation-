# -*- coding: utf-8 -*-
"""
Slayt 5'teki tum shape'leri listeler - konumlari gormek icin
"""
from pathlib import Path
from pptx import Presentation

BASE_DIR  = Path(__file__).parent.parent
PPTX_PATH = BASE_DIR / "BIM320_Alev_Segmentasyonu_Sunum_Final.pptx"

prs = Presentation(str(PPTX_PATH))

slide = prs.slides[4]  # Slayt 5 (0-indexed)
print("=== Slayt 5 - Tum Shapeler ===")
for i, shape in enumerate(slide.shapes):
    name = shape.name
    l = shape.left / 914400
    t = shape.top  / 914400
    w = shape.width / 914400
    h = shape.height / 914400
    txt = ""
    if shape.has_text_frame:
        txt = shape.text_frame.text.strip()[:50].replace("\n", " | ")
    print(f"  [{i}] {name:30s}  L={l:.2f} T={t:.2f} W={w:.2f} H={h:.2f}  |  '{txt}'")
