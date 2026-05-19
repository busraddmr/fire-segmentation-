# -*- coding: utf-8 -*-
"""
BIM320_Alev_Segmentasyonu_Sunum_Final.pptx icindeki
VOTE_THR slaytinda 'GORSEL' placeholder yerine
report_precision_recall.png gorselini ekler.
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

BASE_DIR  = Path(__file__).parent.parent
OUT_DIR   = BASE_DIR / "outputs"
PPTX_PATH = BASE_DIR / "BIM320_Alev_Segmentasyonu_Sunum_Final.pptx"
IMG_PATH  = OUT_DIR / "report_precision_recall.png"

if not IMG_PATH.exists():
    print("HATA: Gorsel bulunamadi:", IMG_PATH)
    sys.exit(1)

prs = Presentation(str(PPTX_PATH))
print("Toplam slayt sayisi:", len(prs.slides))

hedef_slide = None
hedef_idx   = None
left = top = width = height = None

for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip().upper()
            if "GORSEL" in txt or "G\u00D6RSEL" in txt:
                hedef_slide = slide
                hedef_idx   = i
                left   = shape.left
                top    = shape.top
                width  = shape.width
                height = shape.height
                print(f"Placeholder bulundu -> Slayt {i+1}")
                print(f"  Konum: left={left/914400:.2f}in top={top/914400:.2f}in "
                      f"w={width/914400:.2f}in h={height/914400:.2f}in")
                # Placeholder'i sil
                shape._element.getparent().remove(shape._element)
                break
    if hedef_slide:
        break

if hedef_slide is None:
    print("'GORSEL' yazisi bulunamadi. Slayt metinleri:")
    for i, slide in enumerate(prs.slides):
        texts = [s.text_frame.text.strip()[:40]
                 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
        if texts:
            print(f"  Slayt {i+1}: {texts}")
    sys.exit(1)

# Gorseli ekle (kucuk kenar boslugu ile)
MARGIN = Inches(0.08)
hedef_slide.shapes.add_picture(
    str(IMG_PATH),
    left   + MARGIN,
    top    + MARGIN,
    width  - 2 * MARGIN,
    height - 2 * MARGIN
)

print("Gorsel eklendi:", IMG_PATH.name)

prs.save(str(PPTX_PATH))
print("Sunum kaydedildi:", PPTX_PATH.name)
